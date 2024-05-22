from flask import Flask, render_template, request, send_file
import requests
from pptx import Presentation
from pptx.util import Inches
import io
from PIL import Image as PILImage
import openai
import re
from dotenv import load_dotenv
import os

app = Flask(__name__)

# Load environment variables from .env file
load_dotenv()

# Set your OpenAI API key
openai.api_key = os.getenv("OPENAI_API_KEY")

# Set your Pexels API key
PEXELS_API_KEY = os.getenv("PEXELS_API_KEY")

@app.route('/', methods=['GET'])
def index():
    return render_template('index.html')

@app.route('/generate-ppt', methods=['POST'])
def generate_ppt():
    title = request.form['title']
    presenter = request.form['presenter']
    num_slides = int(request.form['num_slides'])
    prompt = request.form['prompt']
    include_references = 'include_references' in request.form
    include_images = 'include_images' in request.form
    template_choice = request.form['template_choice']
    template_path = f'presentations/{template_choice}'

    slides_content, conclusion_content, references_content = generate_content(
        prompt, num_slides, include_references
    )
    ppt_file = create_ppt(slides_content, conclusion_content, references_content, title, presenter, template_path, include_images, num_slides)
    return send_file(ppt_file, as_attachment=True, download_name='presentation.pptx')

def generate_content(prompt, num_slides, include_references):
    response_main = openai.ChatCompletion.create(
        model="gpt-3.5-turbo-0125",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt}
        ],
        max_tokens=200 * num_slides
    )
    main_content = response_main.choices[0].message['content'].strip()
    slides_content = main_content.split("\n\n")[:num_slides]  # Split content into slides based on double newlines

    response_conclusion = openai.ChatCompletion.create(
        model="gpt-3.5-turbo-0125",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": "Summarize the key points covered in the following content: " + main_content}
        ],
        max_tokens=150
    )
    conclusion_content = response_conclusion.choices[0].message['content'].strip()

    references_content = ""
    if include_references:
        response_references = openai.ChatCompletion.create(
            model="gpt-3.5-turbo-0125",
            messages=[
                {"role": "system", "content": "You are a helpful assistant."},
                {"role": "user", "content": "Provide a list of references based on the discussed topics."}
            ],
            max_tokens=150
        )
        references_content = response_references.choices[0].message['content'].strip()

    return slides_content, conclusion_content, references_content

def fetch_image(query):
    headers = {'Authorization': PEXELS_API_KEY}
    params = {'query': query, 'per_page': 1}
    response = requests.get('https://api.pexels.com/v1/search', headers=headers, params=params)
    if response.status_code == 200:
        results = response.json()
        if results['photos']:
            image_url = results['photos'][0]['src']['original']
            image_response = requests.get(image_url)
            if image_response.status_code == 200:
                image = PILImage.open(io.BytesIO(image_response.content))
                return image
    return None

def add_image_to_slide(slide, image, prs):
    image_stream = io.BytesIO()
    image.save(image_stream, format='PNG')
    image_stream.seek(0)
    slide.shapes.add_picture(image_stream, prs.slide_width - Inches(4), Inches(1), width=Inches(3), height=Inches(2))

def create_ppt(slides_content, conclusion_content, references_content, presentation_title, presenter_name, template_path, include_images, num_slides):
    prs = Presentation(template_path)

    # Remove all existing slides from the presentation
    while len(prs.slides) > 0:
        xml_slides = prs.slides._sldIdLst
        prs.part.drop_rel(xml_slides[0].rId)
        del xml_slides[0]

    # Create the title slide
    title_slide_layout = prs.slide_layouts[0]  # Assuming the first layout is for the title
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = presentation_title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Presented by " + presenter_name

    # Create content slides using the specified content layout
    content_slide_layout = prs.slide_layouts[2]  # TITLE_AND_BODY layout
    for content in slides_content[:num_slides]:
        slide = prs.slides.add_slide(content_slide_layout)
        parts = content.split(':', 1)
        title = re.sub(r'\*\*|__|```', '', parts[0]).strip()  # Clean markdown syntax
        body = parts[1].strip() if len(parts) > 1 else ""

        if len(title.split()) > 5:
            body = title + " " + body  # Long title, move to body
            title = "Overview"

        slide.shapes.title.text = title
        slide.placeholders[1].text = body  # Assumed body placeholder exists in this layout

        if include_images and len(title.split()) <= 5:
            image = fetch_image(title)
            if image:
                add_image_to_slide(slide, image, prs)

    # Optional conclusion slide
    if conclusion_content:
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = "Conclusion"
        slide.placeholders[1].text = conclusion_content

    # Optional references slide
    if references_content:
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = "References"
        slide.placeholders[1].text = references_content

    file_path = 'generated_presentation_using_template.pptx'
    prs.save(file_path)
    return file_path

if __name__ == '__main__':
    app.run(debug=True)
